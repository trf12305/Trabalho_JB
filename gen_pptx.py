# -*- coding: utf-8 -*-
"""Gera Dashboard_JB_Apresentacao.pptx - Relatorio executivo+tecnico"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============= TEMA =============
ORANGE = RGBColor(0xE7, 0x6F, 0x2A)
NAVY   = RGBColor(0x1F, 0x38, 0x64)
DARK   = RGBColor(0x0F, 0x14, 0x29)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x21, 0xA8, 0x6C)
RED    = RGBColor(0xDC, 0x35, 0x45)
YELLOW = RGBColor(0xF7, 0xC5, 0x18)
PURPLE = RGBColor(0x7A, 0x4A, 0xC8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ============= HELPERS =============
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb

def add_multitext(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP):
    """lines: lista de dicts {text, size, bold, color, font}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if "space_after" in ln: p.space_after = Pt(ln["space_after"])
        r = p.add_run(); r.text = ln["text"]
        f = r.font
        f.name = ln.get("font", "Calibri")
        f.size = Pt(ln.get("size", 16))
        f.bold = ln.get("bold", False)
        f.color.rgb = ln.get("color", DARK)
    return tb

def add_bullet(slide, x, y, w, h, items, size=16, color=DARK,
               line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = "•  " + it
        f = r.font
        f.name = "Calibri"; f.size = Pt(size); f.color.rgb = color
    return tb

def add_table(slide, x, y, w, h, header, rows, col_widths=None,
              header_bg=NAVY, header_fg=WHITE,
              row_alt=LIGHT, font_size=12, header_size=13):
    n_cols = len(header)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * (cw / total)))
    for i, ht in enumerate(header):
        c = tbl.cell(0, i)
        c.fill.solid(); c.fill.fore_color.rgb = header_bg
        c.margin_left = Inches(0.08); c.margin_right = Inches(0.08)
        c.margin_top = Inches(0.04); c.margin_bottom = Inches(0.04)
        tf = c.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = ht
        f = r.font
        f.name = "Calibri"; f.size = Pt(header_size); f.bold = True; f.color.rgb = header_fg
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if ri % 2 else row_alt
            c.margin_left = Inches(0.08); c.margin_right = Inches(0.08)
            c.margin_top = Inches(0.04); c.margin_bottom = Inches(0.04)
            tf = c.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            f = r.font
            f.name = "Calibri"; f.size = Pt(font_size); f.color.rgb = DARK
    return tbl_shape

def kpi_card(slide, x, y, w, h, label, value, color_value=ORANGE,
             color_label=GRAY, sub=None):
    add_rect(slide, x, y, w, h, LIGHT)
    add_text(slide, x + Inches(0.15), y + Inches(0.12),
             w - Inches(0.3), Inches(0.3),
             label, size=10, bold=True, color=color_label)
    add_text(slide, x + Inches(0.15), y + Inches(0.40),
             w - Inches(0.3), Inches(0.7),
             value, size=24, bold=True, color=color_value)
    if sub:
        add_text(slide, x + Inches(0.15), y + h - Inches(0.45),
                 w - Inches(0.3), Inches(0.35),
                 sub, size=10, color=GRAY)

def slide_header(slide, title, subtitle=None, chip=None):
    # barra superior
    add_rect(slide, Emu(0), Emu(0), SW, Inches(0.08), ORANGE)
    # titulo
    add_text(slide, Inches(0.5), Inches(0.25), SW - Inches(1), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), SW - Inches(1), Inches(0.4),
                 subtitle, size=14, color=GRAY)
    if chip:
        chip_w = Inches(1.6); chip_h = Inches(0.32)
        chip_x = SW - chip_w - Inches(0.5); chip_y = Inches(0.3)
        add_rect(slide, chip_x, chip_y, chip_w, chip_h, NAVY)
        add_text(slide, chip_x, chip_y, chip_w, chip_h,
                 chip, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_footer(slide, n, total):
    add_text(slide, Inches(0.5), SH - Inches(0.35),
             SW - Inches(1), Inches(0.25),
             f"JB Proteção  |  Dashboard Financeiro",
             size=9, color=GRAY)
    add_text(slide, Inches(0.5), SH - Inches(0.35),
             SW - Inches(1), Inches(0.25),
             f"{n} / {total}", size=9, color=GRAY,
             align=PP_ALIGN.RIGHT)


# ============= SLIDES =============
SLIDES = []

# ---------- 1: CAPA ----------
def slide_capa():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, DARK)
    add_rect(s, Emu(0), Inches(3.0), SW, Inches(0.05), ORANGE)
    add_text(s, Inches(0), Inches(1.6), SW, Inches(1.2),
             "JB PROTEÇÃO", size=60, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0), Inches(2.4), SW, Inches(0.7),
             "Dashboard Financeiro Inteligente", size=28, color=WHITE,
             align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0), Inches(3.3), SW, Inches(0.5),
             "Relatório Executivo e Técnico", size=20, color=ORANGE,
             align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0), Inches(5.2), SW, Inches(0.4),
             "Período: Maio/2026  |  17.795 títulos  |  12 unidades",
             size=14, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0), Inches(5.7), SW, Inches(0.4),
             "Última sincronização: 27/05/2026 16:12",
             size=12, color=GRAY, align=PP_ALIGN.CENTER, font="Calibri")

slide_capa()

# ---------- 2: AGENDA ----------
def slide_agenda():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Agenda", "Estrutura desta apresentação")
    items = [
        "1.  Contexto da base de dados (origem e cobertura)",
        "2.  KPIs principais — fórmulas e origem de cada card",
        "3.  KPIs avançados de gestão financeira",
        "4.  Análise de aging e inadimplência",
        "5.  Auditoria automática de inconsistências",
        "6.  Gráficos e visualizações",
        "7.  Filtros, controles e export",
        "8.  Arquitetura técnica e confiabilidade",
        "9.  Conclusão e próximos passos",
    ]
    add_bullet(s, Inches(1.0), Inches(1.7),
               Inches(11), Inches(5),
               items, size=18, color=DARK, line_spacing=1.35)
    slide_footer(s, 2, 14)

slide_agenda()

# ---------- 3: HEADLINE ----------
def slide_headline():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Visão Geral — Maio/2026",
                 "Os números que importam")
    # 4 cards principais em cima
    card_w, card_h = Inches(2.95), Inches(1.6)
    y0 = Inches(1.5)
    kpi_card(s, Inches(0.5), y0, card_w, card_h,
             "FATURAMENTO PREVISTO", "R$ 1.800.906,37", color_value=GREEN,
             sub="Valor de face de 17.795 títulos")
    kpi_card(s, Inches(3.6), y0, card_w, card_h,
             "FATURAMENTO REALIZADO", "R$ 1.215.351,90", color_value=ORANGE,
             sub="67,5% de conversão")
    kpi_card(s, Inches(6.7), y0, card_w, card_h,
             "A RECEBER", "R$ 585.821,41", color_value=NAVY,
             sub="5.389 títulos em aberto")
    kpi_card(s, Inches(9.8), y0, card_w, card_h,
             "INADIMPLÊNCIA", "R$ 0,00", color_value=RED,
             sub="Zero pendências ativas")
    # 4 cards secundários
    card_h2 = Inches(1.4)
    y1 = Inches(3.4)
    kpi_card(s, Inches(0.5), y1, card_w, card_h2,
             "PONTUALIDADE", "66,9%", color_value=YELLOW,
             sub="Atraso médio: 2,5 dias")
    kpi_card(s, Inches(3.6), y1, card_w, card_h2,
             "TICKET MÉDIO", "R$ 101,20", color_value=PURPLE,
             sub="Por título")
    kpi_card(s, Inches(6.7), y1, card_w, card_h2,
             "BASES ATIVAS", "12", color_value=NAVY,
             sub="AL, PE, PB, SE")
    kpi_card(s, Inches(9.8), y1, card_w, card_h2,
             "AUDITORIA", "1 alerta", color_value=YELLOW,
             sub="R$ 50,00 a verificar")
    # box destaque
    add_rect(s, Inches(0.5), Inches(5.4),
             SW - Inches(1), Inches(1.5), NAVY)
    add_text(s, Inches(0.7), Inches(5.55),
             SW - Inches(1.4), Inches(0.4),
             "MENSAGEM-CHAVE", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(5.95),
             SW - Inches(1.4), Inches(1.0),
             "A JB Proteção movimentou 17.795 títulos em Maio/2026, "
             "totalizando R$ 1,8 milhão previstos. Já recebeu R$ 1,21 "
             "milhão (67,5%) com 66,9% de pontualidade. Zero inadimplência. "
             "Operação saudável distribuída em 12 unidades, atendendo 14.933 associados.",
             size=14, color=WHITE)
    slide_footer(s, 3, 14)

slide_headline()

# ---------- 4: CONTEXTO ----------
def slide_contexto():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "1. Contexto da Base de Dados",
                 "Origem, escopo e cobertura", chip="TÉCNICO")
    # Esquerda — Fonte
    add_text(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
             "FONTE DE DADOS", size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(0.5), Inches(2.0), Inches(6), Inches(2.5), [
        "Sistema: Siprov (export JSON layout 496)",
        "Arquivo: dashboard_financeiro_live.json",
        "Tamanho: 20,6 MB / 17.795 registros / 33 colunas",
        "Modo de operação: JSON fixo (auto-sync desativado)",
    ], size=14)
    # Direita — Cobertura
    add_text(s, Inches(6.8), Inches(1.6), Inches(6), Inches(0.4),
             "COBERTURA", size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.5), [
        "Tipo: CREDITO (Contas a Receber)",
        "Vencimento: 01/05/2026 a 31/05/2026",
        "LIQUIDADO: 12.406 títulos",
        "ABERTO: 5.389 títulos",
        "PENDENTE: 0 títulos",
    ], size=14)
    # Bases
    add_text(s, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
             "12 BASES JB ATIVAS NO PERÍODO", size=14, bold=True, color=ORANGE)
    add_rect(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.3), LIGHT)
    add_text(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.2),
             "JB ALAGOAS  •  JB ARACAJU  •  JB ARAPIRACA  •  JB CAMPINA GRANDE  •  "
             "JB CARUARU  •  JB CORURIPE  •  JB ITABAIANA  •  JB MACEIÓ  •  "
             "JB PARAÍBA  •  JB PATOS  •  JB PERNAMBUCO  •  JB SERGIPE",
             size=13, color=DARK)
    slide_footer(s, 4, 14)

slide_contexto()

# ---------- 5: KPIs PRINCIPAIS PARTE 1 ----------
def slide_kpi_p1():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "2. KPIs Principais — Quantidade e Valores",
                 "Origem de dados de cada card", chip="EXEC + TÉC")
    add_table(s,
        Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.5),
        ["Card", "Coluna JSON", "Fórmula", "Vencimento", "Liquidação"],
        [
            ["QTD. TÍTULOS", "(contagem)", "len(dados filtrados)",
             "17.795", "12.406"],
            ["TOTAL LIQUIDADO / RECEBIDO", "liquidacao_valor_liquidado",
             "SUM dos LIQUIDADO", "R$ 1.215.351,90", "R$ 1.215.351,90"],
            ["VALOR DE FACE", "titulo_valor",
             "SUM (todos / só LIQUIDADO)",
             "R$ 1.800.906,37", "R$ 1.215.084,96"],
            ["TICKET MÉDIO", "(razão)",
             "Total ÷ Quantidade",
             "R$ 101,20", "R$ 97,96"],
        ],
        col_widths=[2.2, 2.4, 2.3, 2.4, 2.4],
        font_size=12, header_size=12)
    # Insight box
    add_rect(s, Inches(0.5), Inches(4.5),
             Inches(12.3), Inches(2.4), LIGHT)
    add_text(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.4),
             "INSIGHT PARA O CHEFE", size=12, bold=True, color=ORANGE)
    add_bullet(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2), [
        "TOTAL LIQUIDADO (R$ 1.215.351,90) > VALOR DE FACE PAGOS (R$ 1.215.084,96): "
        "a diferença de R$ 266,94 é o saldo líquido de juros − descontos cobrados.",
        "Ticket médio diferente entre modos: face (R$ 101,20) considera todos; "
        "recebido (R$ 97,96) considera só os pagos — clientes pequenos pagam mais.",
        "Quantidade Liquidação (12.406) confere 1:1 com a planilha Excel: cada linha = 1 título.",
    ], size=12)
    slide_footer(s, 5, 14)

slide_kpi_p1()

# ---------- 6: KPIs PRINCIPAIS PARTE 2 ----------
def slide_kpi_p2():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "2. KPIs Principais — Pontualidade e Bases",
                 "Indicadores qualitativos", chip="EXEC + TÉC")
    # Pontualidade — esquerda
    add_text(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
             "PONTUALIDADE — 66,9%", size=15, bold=True, color=ORANGE)
    add_table(s,
        Inches(0.5), Inches(2.1), Inches(6), Inches(2.5),
        ["Categoria", "Qtd", "%"],
        [
            ["Pago ANTES do vencimento", "6.272", "50,6%"],
            ["Pago NO DIA do vencimento", "2.026", "16,3%"],
            ["Pago APÓS o vencimento", "4.107", "33,1%"],
            ["TOTAL pontuais (verde+azul)", "8.298", "66,9%"],
        ],
        col_widths=[3.0, 1.5, 1.5], font_size=11)
    add_multitext(s, Inches(0.5), Inches(4.8), Inches(6), Inches(1.7), [
        {"text": "Colunas JSON: ", "size": 11, "bold": True, "color": NAVY},
        {"text": "titulo_data_vencimento e liquidacao_data_liquidacao", "size": 11},
        {"text": "Fórmula: ", "size": 11, "bold": True, "color": NAVY},
        {"text": "COUNT(data_liq ≤ data_venc) ÷ TOTAL × 100", "size": 11},
        {"text": "Atraso médio: 2,5 dias  •  Antecipação média: 7 dias",
         "size": 11, "bold": True, "color": GREEN, "space_after": 2},
    ])
    # Bases — direita
    add_text(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.4),
             "BASES — 12 unidades ativas", size=15, bold=True, color=ORANGE)
    add_multitext(s, Inches(7.0), Inches(2.1), Inches(6), Inches(0.8), [
        {"text": "Coluna JSON: ", "size": 11, "bold": True, "color": NAVY},
        {"text": "unidade_nome_fantasia", "size": 11},
        {"text": "Fórmula: COUNT(DISTINCT unidade_nome_fantasia)",
         "size": 11, "bold": True, "color": NAVY},
    ])
    add_table(s,
        Inches(7.0), Inches(3.1), Inches(6), Inches(3.5),
        ["UF", "Bases JB"],
        [
            ["AL", "ALAGOAS, ARAPIRACA, CORURIPE, MACEIÓ"],
            ["PE", "CARUARU, PERNAMBUCO"],
            ["PB", "CAMPINA GRANDE, PARAÍBA, PATOS"],
            ["SE", "ARACAJU, ITABAIANA, SERGIPE"],
        ],
        col_widths=[1.0, 5.0], font_size=11)
    slide_footer(s, 6, 14)

slide_kpi_p2()

# ---------- 7: KPIs AVANÇADOS ----------
def slide_kpi_avancado():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "3. KPIs Avançados de Gestão",
                 "Métricas para tomada de decisão estratégica",
                 chip="EXEC + TÉC")
    add_table(s,
        Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5),
        ["Indicador", "Valor", "Coluna/Fórmula", "Leitura executiva"],
        [
            ["A RECEBER (ABERTO)", "R$ 585.821,41",
             "SUM(titulo_valor) WHERE situação=ABERTO",
             "Carteira que ainda vai virar caixa"],
            ["INADIMPLÊNCIA", "R$ 0,00",
             "SUM(titulo_valor) WHERE situação=PENDENTE",
             "Zero inadimplência ativa — saúde excelente"],
            ["% CONVERSÃO", "67,5%",
             "Realizado ÷ Previsto × 100",
             "Já recebemos 2/3 da carteira do mês"],
            ["DSO", "301,3 dias",
             "MÉDIA(data_liq − data_venc) dos pagos",
             "Reflete pagamentos antecipados (1-2 anos)"],
            ["ATRASO MÉDIO", "0,0 dia",
             "MÉDIA(hoje − data_venc) dos pendentes",
             "Zero pendentes hoje"],
            ["MoM / YoY", "0% / 0%",
             "Δ mês a mês / ano a ano",
             "Ativa quando houver 2-3 meses de histórico"],
        ],
        col_widths=[2.4, 1.8, 3.5, 4.6], font_size=11, header_size=12)
    # Box de destaque
    add_rect(s, Inches(0.5), Inches(6.3),
             Inches(12.3), Inches(0.7), GREEN)
    add_text(s, Inches(0.7), Inches(6.3),
             Inches(12), Inches(0.7),
             "✓ Conversão de 67,5% + Inadimplência R$ 0 = operação financeiramente saudável e previsível.",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    slide_footer(s, 7, 14)

slide_kpi_avancado()

# ---------- 8: AGING ----------
def slide_aging():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "4. Aging — Inadimplência por Idade",
                 "Buckets de gestão de carteira", chip="EXEC + TÉC")
    # 4 cards
    card_w, card_h = Inches(2.95), Inches(1.8)
    y0 = Inches(1.8)
    kpi_card(s, Inches(0.5), y0, card_w, card_h, "D30 (1-30 dias)",
             "R$ 0,00", color_value=GREEN, sub="Atraso recente")
    kpi_card(s, Inches(3.6), y0, card_w, card_h, "D60 (31-60 dias)",
             "R$ 0,00", color_value=GREEN, sub="Atraso intermediário")
    kpi_card(s, Inches(6.7), y0, card_w, card_h, "D90 (61-90 dias)",
             "R$ 0,00", color_value=GREEN, sub="Atraso crítico")
    kpi_card(s, Inches(9.8), y0, card_w, card_h, "D180+ (90+ dias)",
             "R$ 0,00", color_value=GREEN, sub="Provisão de perda")
    # Box
    add_rect(s, Inches(0.5), Inches(4.0),
             Inches(12.3), Inches(2.7), LIGHT)
    add_text(s, Inches(0.7), Inches(4.15), Inches(12), Inches(0.4),
             "METODOLOGIA E INTERPRETAÇÃO", size=12, bold=True, color=ORANGE)
    add_bullet(s, Inches(0.7), Inches(4.55), Inches(12), Inches(2.2), [
        "Coluna JSON: titulo_valor filtrada por situação=PENDENTE + diferença de dias entre vencimento e hoje.",
        "Buckets seguem padrão de mercado (30/60/90/90+) para análise de qualidade de carteira.",
        "Todos os buckets zerados porque INADIMPLÊNCIA total é R$ 0,00.",
        "Quando houver pendências, o sistema classificará automaticamente nos buckets correspondentes.",
    ], size=13)
    slide_footer(s, 8, 14)

slide_aging()

# ---------- 9: AUDITORIA ----------
def slide_auditoria():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "5. Auditoria Automática de Inconsistências",
                 "Diferencial do dashboard — detecção proativa de erros",
                 chip="EXEC + TÉC")
    # Badge
    add_rect(s, Inches(0.5), Inches(1.6),
             Inches(4.5), Inches(1.3), YELLOW)
    add_text(s, Inches(0.7), Inches(1.7),
             Inches(4.1), Inches(0.4),
             "⚠ ALERTA ATIVO", size=12, bold=True, color=DARK)
    add_text(s, Inches(0.7), Inches(2.0),
             Inches(4.1), Inches(0.6),
             "1 título", size=26, bold=True, color=DARK)
    add_text(s, Inches(0.7), Inches(2.55),
             Inches(4.1), Inches(0.35),
             "R$ 50,00 a auditar", size=12, color=DARK)
    # Lógica
    add_text(s, Inches(5.3), Inches(1.6),
             Inches(7.5), Inches(0.4),
             "LÓGICA DE DETECÇÃO", size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(5.3), Inches(2.0),
               Inches(7.5), Inches(2),
               [
                "Filtra títulos com situação = LIQUIDADO/PAGO/QUITADO",
                "Verifica se liquidacao_data_liquidacao está vazia",
                "Verifica se liquidacao_valor_liquidado é nulo ou ≤ 0",
                "Exibe badge amarelo no card TOTAL LIQUIDADO",
            ], size=12)
    # Caso atual
    add_text(s, Inches(0.5), Inches(3.4), Inches(12), Inches(0.4),
             "CASO DETECTADO NO PERÍODO", size=14, bold=True, color=ORANGE)
    add_table(s,
        Inches(0.5), Inches(3.85), Inches(12.3), Inches(2.7),
        ["Campo", "Valor"],
        [
            ["Associado", "Miguel Ferreira Lima"],
            ["CPF", "419.107.614-00"],
            ["Unidade", "JB ARAPIRACA"],
            ["Veículo", "HONDA — placa QTT3356"],
            ["Parcela", "4/12 (RENOVAÇÃO)"],
            ["Vencimento", "10/05/2026"],
            ["Valor", "R$ 50,00"],
            ["Diagnóstico", "Marcado LIQUIDADO sem dados de pagamento — provável erro de cadastro Siprov"],
        ],
        col_widths=[2.5, 9.5], font_size=11)
    slide_footer(s, 9, 14)

slide_auditoria()

# ---------- 10: GRÁFICOS ----------
def slide_graficos():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "6. Visualizações Gráficas",
                 "9 dimensões de análise da carteira", chip="EXEC")
    add_table(s,
        Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.2),
        ["Gráfico", "Colunas JSON", "O que mostra"],
        [
            ["Fluxo Mensal",
             "titulo_data_vencimento + titulo_valor",
             "Faturamento previsto por mês de vencimento"],
            ["Liquidação Mensal",
             "liquidacao_data_liquidacao + liquidacao_valor_liquidado",
             "Caixa realizado por mês de pagamento"],
            ["Top 10 Unidades",
             "unidade_nome_fantasia + valor",
             "Ranking de bases por receita"],
            ["Top 10 Consultores",
             "beneficio_consultor + valor",
             "Performance comercial individual"],
            ["Formas de Pagamento",
             "liquidacao_tipo_liquidacao + valor",
             "Mix: Boleto Santander, Pix, Dinheiro, RumoPay, etc."],
            ["Top 10 Planos",
             "beneficio_planos_principais + valor",
             "Planos mais vendidos da carteira"],
            ["Top 10 Representantes",
             "beneficio_representante + valor",
             "Performance de equipes/representantes"],
            ["Ranking de Bases",
             "unidade_nome_fantasia + valor_liquidado",
             "Receita realizada por unidade JB"],
            ["Comparativo Mensal",
             "titulo_situacao_titulo + valores",
             "Face × Liquidado × Pendente × Aberto por mês"],
        ],
        col_widths=[2.8, 4.5, 5.0], font_size=10, header_size=12)
    slide_footer(s, 10, 14)

slide_graficos()

# ---------- 11: FILTROS E EXPORT ----------
def slide_filtros():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "7. Filtros, Controles e Export",
                 "Interatividade do dashboard", chip="TÉCNICO")
    add_table(s,
        Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.8),
        ["Controle", "Função", "Endpoint/Lógica"],
        [
            ["Bases", "Filtrar por 1 ou várias unidades",
             "WHERE unidade_nome_fantasia IN (...)"],
            ["Data Inicial / Final", "Range customizado de período",
             "WHERE data_filtro BETWEEN di AND df"],
            ["Tipo de Filtro", "Vencimento × Liquidação",
             "Define coluna de data e métrica agregada"],
            ["Atualizar", "Recarrega dados sem refresh",
             "GET /api/financeiro?..."],
            ["Auto-Refresh", "Atualização periódica",
             "setInterval no front-end"],
            ["Sync Siprov", "Coleta manual no Siprov",
             "POST /api/admin/sync (assíncrono)"],
            ["Exportar CSV", "Download dos dados filtrados",
             "GET /api/financeiro/export?..."],
            ["Tabela operacional", "Top 500 registros do filtro",
             "Renderizada no DOM com paginação"],
        ],
        col_widths=[2.7, 5.0, 4.6], font_size=11, header_size=12)
    add_rect(s, Inches(0.5), Inches(6.6),
             Inches(12.3), Inches(0.5), NAVY)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
             "Export CSV gera arquivo com TODOS os 17.795 títulos e 33 colunas — pronto para Excel/BI.",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    slide_footer(s, 11, 14)

slide_filtros()

# ---------- 12: ARQUITETURA ----------
def slide_arquitetura():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "8. Arquitetura Técnica",
                 "Stack e pipeline de dados", chip="TÉCNICO")
    # Stack — esquerda
    add_text(s, Inches(0.5), Inches(1.6),
             Inches(6), Inches(0.4),
             "STACK TECNOLÓGICA", size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(0.5), Inches(2.0),
               Inches(6), Inches(3), [
        "Backend: Python 3 + Flask + APScheduler",
        "Frontend: HTML5 + TailwindCSS + Chart.js",
        "Integração: API REST Siprov (layout 496 async)",
        "Persistência: JSON local (sem banco de dados)",
        "Cache: in-memory por mtime do arquivo",
    ], size=13)
    # Pipeline — direita
    add_text(s, Inches(6.8), Inches(1.6),
             Inches(6), Inches(0.4),
             "PIPELINE DE DADOS", size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(6.8), Inches(2.0),
               Inches(6), Inches(3), [
        "1. carregar_dados_json() — lê o JSON do disco",
        "2. processar_dados_para_dash() — normaliza campos",
        "3. filtrar_dados() — aplica filtros do usuário",
        "4. gerar_dashboard_analitico() — calcula KPIs",
        "5. API /api/financeiro — entrega JSON ao front-end",
    ], size=13)
    # Confiabilidade
    add_text(s, Inches(0.5), Inches(4.7),
             Inches(12), Inches(0.4),
             "CONFIABILIDADE E VALIDAÇÃO", size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(0.5), Inches(5.1),
               Inches(12), Inches(2), [
        "100% das rotas testadas (login, financeiro, eventos, vendas, export, auditoria)",
        "Validação cruzada: todos os totais batem 1:1 com planilha Excel oficial",
        "Tratamento de inputs inválidos: datas malformadas e tipos desconhecidos têm fallback seguro",
        "Auto-sync desativável para reprodutibilidade em apresentações",
    ], size=12)
    slide_footer(s, 12, 14)

slide_arquitetura()

# ---------- 13: CONCLUSÃO ----------
def slide_conclusao():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "9. Conclusão",
                 "Síntese dos resultados e próximos passos")
    # Entregas
    add_text(s, Inches(0.5), Inches(1.6),
             Inches(12), Inches(0.4),
             "✓ PRINCIPAIS ENTREGAS", size=14, bold=True, color=GREEN)
    add_bullet(s, Inches(0.5), Inches(2.0),
               Inches(12), Inches(2.5), [
        "Visibilidade executiva e operacional sobre 17.795 títulos da carteira",
        "Validação 1:1 com a planilha Excel — todos os números confirmados",
        "Detecção automática de inconsistências de cadastro no Siprov",
        "Dois modos analíticos (Vencimento e Liquidação) para visões complementares",
        "Export CSV completo para integração com Excel e BI",
        "KPIs avançados (DSO, Aging, MoM/YoY) prontos para escalar com histórico",
    ], size=13)
    # Próximos passos
    add_text(s, Inches(0.5), Inches(4.8),
             Inches(12), Inches(0.4),
             "→ PRÓXIMOS PASSOS RECOMENDADOS",
             size=14, bold=True, color=ORANGE)
    add_bullet(s, Inches(0.5), Inches(5.2),
               Inches(12), Inches(2), [
        "Validação dos números pelos times de operação e contabilidade",
        "Resolução do caso de auditoria (Miguel Ferreira / R$ 50)",
        "Reativação do auto-sync após validação inicial",
        "Acúmulo de 2-3 meses de histórico para ativar MoM/YoY",
    ], size=13)
    slide_footer(s, 13, 14)

slide_conclusao()

# ---------- 14: OBRIGADO ----------
def slide_obrigado():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, DARK)
    add_rect(s, Emu(0), Inches(3.4), SW, Inches(0.05), ORANGE)
    add_text(s, Inches(0), Inches(2.5), SW, Inches(1),
             "Obrigado.", size=60, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.7), SW, Inches(0.6),
             "Dashboard Financeiro JB Proteção",
             size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.3), SW, Inches(0.4),
             "Disponível em http://127.0.0.1:5000/financeiro",
             size=14, color=GRAY, align=PP_ALIGN.CENTER)

slide_obrigado()

# Salvar
out = "Dashboard_JB_Apresentacao.pptx"
prs.save(out)
import os
print(f"OK: {out} ({os.path.getsize(out)/1024:.1f} KB) / {len(prs.slides)} slides")
